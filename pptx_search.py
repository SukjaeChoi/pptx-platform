"""
pptx_search.py
──────────────
지정한 디렉토리(및 하위 디렉토리)의 모든 .pptx 파일에서 키워드를 검색하여
결과를 엑셀 파일로 저장합니다.

사용법:
    python pptx_search.py <디렉토리_경로> <검색어> [출력파일.xlsx]

예시:
    python pptx_search.py "C:\\강의자료" "파인튜닝"
    python pptx_search.py "C:\\강의자료" "LLM" 결과.xlsx

검색 범위:
    - 슬라이드 본문 텍스트 (모든 도형 · 텍스트박스 · 표)
    - 슬라이드 노트(발표자 노트)
    - 대소문자 구분 없이 검색 (한글 포함)
"""

import sys
import os
import re
from pathlib import Path
from datetime import datetime

# ── 패키지 확인 ───────────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
    from pptx.util import Pt
except ImportError:
    print("필수 패키지가 없습니다:\n  pip install python-pptx openpyxl")
    sys.exit(1)

try:
    from openpyxl import Workbook
    from openpyxl.styles import (Font, PatternFill, Alignment,
                                  Border, Side, GradientFill)
    from openpyxl.utils import get_column_letter
except ImportError:
    print("필수 패키지가 없습니다:\n  pip install openpyxl")
    sys.exit(1)


# ── 유틸리티 ──────────────────────────────────────────────────────────────────

def sanitize(text: str) -> str:
    """openpyxl 불허 제어 문자 제거 및 공백 정리."""
    if not isinstance(text, str):
        return text
    text = text.replace("\r\n", " ").replace("\r", " ").replace("\n", " ").replace("\t", " ")
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", text)
    return re.sub(r" +", " ", text).strip()


def safe_write(ws, row, col, value, **kwargs):
    """셀에 값을 안전하게 씁니다 (문자열 자동 정제)."""
    if isinstance(value, str):
        value = sanitize(value)
    cell = ws.cell(row=row, column=col, value=value)
    for attr, val in kwargs.items():
        setattr(cell, attr, val)
    return cell


# ── PPTX 텍스트 추출 ──────────────────────────────────────────────────────────

def _extract_from_shape(shape, items: list) -> None:
    """단일 도형에서 텍스트를 재귀적으로 추출합니다 (그룹·SmartArt 포함)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        shape_type = shape.shape_type
    except Exception:
        shape_type = None

    # ① 그룹 도형 → 재귀 탐색
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for child in shape.shapes:
                _extract_from_shape(child, items)
        except Exception:
            pass
        return

    # ② 일반 텍스트 프레임 (텍스트박스, 플레이스홀더, 도형 등)
    if shape.has_text_frame:
        try:
            text = shape.text_frame.text
            if text.strip():
                items.append(("본문", text))
        except Exception:
            pass

    # ③ 표(Table)
    if shape.has_table:
        try:
            for row in shape.table.rows:
                for cell in row.cells:
                    try:
                        text = cell.text_frame.text
                        if text.strip():
                            items.append(("표", text))
                    except Exception:
                        pass
        except Exception:
            pass

    # ④ SmartArt / 기타 — XML에서 직접 텍스트 추출 (fallback)
    if not shape.has_text_frame and not shape.has_table:
        try:
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            raw_texts = shape._element.findall(".//a:t", ns)
            combined = " ".join(t.text for t in raw_texts if t.text and t.text.strip())
            if combined.strip():
                items.append(("기타(XML)", combined))
        except Exception:
            pass


def iter_slide_texts(slide) -> list[tuple[str, str]]:
    """
    슬라이드에서 (출처, 텍스트) 쌍을 모두 반환합니다.
    그룹 도형 내부 및 SmartArt XML까지 재귀 탐색합니다.
    출처: '본문' | '표' | '기타(XML)' | '노트'
    """
    items = []

    for shape in slide.shapes:
        _extract_from_shape(shape, items)

    # 발표자 노트
    if slide.has_notes_slide:
        try:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf is not None:
                note_text = notes_tf.text.strip()
                if note_text:
                    items.append(("노트", note_text))
        except Exception:
            pass

    return items


def get_slide_title(slide) -> str:
    """슬라이드 제목 플레이스홀더 텍스트를 반환합니다."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in (0, 1):
            try:
                text = ph.text_frame.text.strip()
                if text:
                    return sanitize(text)
            except Exception:
                pass
    return "(제목 없음)"


def extract_context(text: str, keyword: str, window: int = 80) -> str:
    """
    텍스트에서 키워드 주변 문자열을 추출합니다 (미리보기용).
    최대 5개 위치까지 표시.
    """
    pattern = re.compile(re.escape(keyword), re.IGNORECASE)
    matches = list(pattern.finditer(text))
    if not matches:
        return ""

    snippets = []
    for m in matches[:5]:
        start = max(0, m.start() - window)
        end   = min(len(text), m.end() + window)
        snippet = text[start:end].replace("\n", " ").replace("\r", " ")
        # 앞뒤 잘림 표시
        if start > 0:
            snippet = "…" + snippet
        if end < len(text):
            snippet = snippet + "…"
        snippets.append(snippet)

    result = "  /  ".join(snippets)
    if len(matches) > 5:
        result += f"  (외 {len(matches)-5}건)"
    return result


# ── 검색 핵심 로직 ────────────────────────────────────────────────────────────

def search_pptx_files(root_dir: Path, keywords: list[str]) -> list[dict]:
    """
    root_dir 이하 모든 .pptx 파일에서 keywords(복수)를 검색합니다.
    각 키워드를 OR 조건으로 검색하며, 매칭된 키워드를 결과에 포함합니다.

    반환 항목 (히트마다 1행):
        파일명, 상대경로, 절대경로, 슬라이드번호, 슬라이드제목,
        출처(본문/표/노트), 매칭키워드, 키워드별횟수, 문맥미리보기
    """
    results = []
    patterns = [(kw, re.compile(re.escape(kw), re.IGNORECASE)) for kw in keywords]
    pptx_files = sorted(root_dir.rglob("*.pptx"))

    if not pptx_files:
        print(f"⚠  '{root_dir}' 에서 .pptx 파일을 찾지 못했습니다.")
        return results

    total = len(pptx_files)
    hit_files = 0
    kw_label = ", ".join(f"[{k}]" for k in keywords)

    print(f"검색어: {kw_label}  |  탐색 파일 수: {total}개\n{'─'*60}")

    for idx, filepath in enumerate(pptx_files, 1):
        rel_path = filepath.relative_to(root_dir)
        print(f"[{idx:>3}/{total}] {rel_path}", end=" ... ", flush=True)

        try:
            prs = Presentation(filepath)
        except PackageNotFoundError:
            print("⚠ 파일 손상 (건너뜀)")
            continue
        except Exception as e:
            print(f"⚠ 읽기 오류: {e} (건너뜀)")
            continue

        file_hits = 0

        for slide_no, slide in enumerate(prs.slides, 1):
            slide_title = get_slide_title(slide)

            for source, text in iter_slide_texts(slide):
                # 각 키워드별 등장 횟수 계산
                kw_counts = {kw: len(pat.findall(text)) for kw, pat in patterns}
                matched = {kw: cnt for kw, cnt in kw_counts.items() if cnt > 0}
                if not matched:
                    continue

                total_count = sum(matched.values())
                matched_kws = ", ".join(matched.keys())

                # 문맥: 첫 번째 매칭 키워드 기준으로 추출
                first_kw = next(iter(matched))
                context = extract_context(text, first_kw)
                # 나머지 키워드도 문맥에 포함
                for kw in list(matched.keys())[1:]:
                    extra = extract_context(text, kw)
                    if extra and extra not in context:
                        context += "  /  " + extra

                results.append({
                    "파일명":       filepath.name,
                    "상대 경로":    str(rel_path),
                    "절대 경로":    str(filepath),
                    "슬라이드 번호": slide_no,
                    "슬라이드 제목": slide_title,
                    "출처":         source,
                    "매칭 키워드":  matched_kws,
                    "등장 횟수":    total_count,
                    "문맥 미리보기": context,
                })
                file_hits += total_count

        if file_hits:
            hit_files += 1
            print(f"✅ {file_hits}건 발견")
        else:
            print("없음")

    print(f"\n{'─'*60}")
    print(f"검색 완료 — 총 {len(results)}건 / {hit_files}개 파일에서 발견")
    return results


# ── 엑셀 저장 ─────────────────────────────────────────────────────────────────

# 스타일 상수
C_DARK   = "1F3864"   # 헤더 배경 (진파랑)
C_MID    = "2E75B6"   # 소계 배경 (중간파랑)
C_LIGHT  = "EBF0FA"   # 짝수행 배경 (연파랑)
C_HIT    = "FFF2CC"   # 검색어 등장 강조 (연노랑)
C_WHITE  = "FFFFFF"
C_RED    = "C00000"

HDR_FONT  = Font(name="Arial", bold=True, color=C_WHITE, size=10)
DATA_FONT = Font(name="Arial", size=10)
SUB_FONT  = Font(name="Arial", bold=True, size=10, color=C_WHITE)
CENTER    = Alignment(horizontal="center", vertical="center", wrap_text=True)
LEFT      = Alignment(horizontal="left",   vertical="center", wrap_text=True)

def thin_border(color="D0D0D0"):
    s = Side(style="thin", color=color)
    return Border(left=s, right=s, top=s, bottom=s)

def fill(color):
    return PatternFill("solid", start_color=color)


def save_to_excel(results: list[dict], keyword: str, output_path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "검색 결과"

    # ── 헤더 타이틀 ──────────────────────────────────────────────────────────
    ws.merge_cells("A1:H1")
    c = ws["A1"]
    c.value     = f'키워드 검색 결과: "{keyword}"'
    c.font      = Font(name="Arial", bold=True, size=14, color=C_DARK)
    c.alignment = CENTER
    ws.row_dimensions[1].height = 28

    ws.merge_cells("A2:H2")
    c = ws["A2"]
    c.value     = (f"생성일시: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
                   f"  |  총 {len(results)}건 발견")
    c.font      = Font(name="Arial", size=9, color="666666")
    c.alignment = CENTER
    ws.row_dimensions[2].height = 16

    # ── 컬럼 헤더 ────────────────────────────────────────────────────────────
    HEADERS = ["번호", "파일명", "상대 경로", "슬라이드\n번호",
               "슬라이드 제목", "출처", "매칭 키워드", "등장\n횟수", "문맥 미리보기"]
    WIDTHS  = [7, 28, 42, 10, 30, 8, 20, 8, 100]

    for col, (h, w) in enumerate(zip(HEADERS, WIDTHS), 1):
        c = ws.cell(row=3, column=col, value=h)
        c.font      = HDR_FONT
        c.fill      = fill(C_DARK)
        c.alignment = CENTER
        c.border    = thin_border()
        ws.column_dimensions[get_column_letter(col)].width = w

    ws.row_dimensions[3].height = 28
    ws.freeze_panes = "A4"

    # ── 데이터 행 ────────────────────────────────────────────────────────────
    prev_file = None
    file_color_toggle = False   # 파일 바뀔 때마다 배경색 교체

    for row_no, rec in enumerate(results, 4):
        # 파일이 바뀌면 색 전환
        if rec["파일명"] != prev_file:
            prev_file = rec["파일명"]
            file_color_toggle = not file_color_toggle

        row_fill = fill(C_LIGHT) if file_color_toggle else fill(C_WHITE)
        bdr      = thin_border()

        values = [
            row_no - 3,
            rec["파일명"],
            rec["상대 경로"],
            rec["슬라이드 번호"],
            rec["슬라이드 제목"],
            rec["출처"],
            rec["매칭 키워드"],
            rec["등장 횟수"],
            rec["문맥 미리보기"],
        ]
        aligns = [CENTER, LEFT, LEFT, CENTER, LEFT, CENTER, CENTER, CENTER, LEFT]

        for col, (val, aln) in enumerate(zip(values, aligns), 1):
            c = ws.cell(row=row_no, column=col,
                        value=sanitize(val) if isinstance(val, str) else val)
            c.font      = DATA_FONT
            c.alignment = aln
            c.border    = bdr
            c.fill      = row_fill

        ws.row_dimensions[row_no].height = 20

    # ── 파일별 소계 시트 ─────────────────────────────────────────────────────
    ws2 = wb.create_sheet("파일별 요약")

    ws2.merge_cells("A1:E1")
    c = ws2["A1"]
    c.value     = f'키워드 "{keyword}" 파일별 요약'
    c.font      = Font(name="Arial", bold=True, size=13, color=C_DARK)
    c.alignment = CENTER
    ws2.row_dimensions[1].height = 26

    for col, (h, w) in enumerate(
        zip(["번호", "파일명", "상대 경로", "히트 슬라이드 수", "총 등장 횟수"],
            [7, 28, 42, 18, 14]), 1):
        c = ws2.cell(row=2, column=col, value=h)
        c.font = HDR_FONT; c.fill = fill(C_DARK)
        c.alignment = CENTER; c.border = thin_border()
        ws2.column_dimensions[get_column_letter(col)].width = w
    ws2.row_dimensions[2].height = 22
    ws2.freeze_panes = "A3"

    # 파일별 집계
    from collections import defaultdict
    file_stats: dict[str, dict] = {}
    for rec in results:
        key = rec["절대 경로"]
        if key not in file_stats:
            file_stats[key] = {
                "파일명": rec["파일명"],
                "상대 경로": rec["상대 경로"],
                "슬라이드 수": set(),
                "총 횟수": 0,
            }
        file_stats[key]["슬라이드 수"].add(rec["슬라이드 번호"])
        file_stats[key]["총 횟수"] += rec["등장 횟수"]

    for i, (_, stat) in enumerate(file_stats.items(), 1):
        row = i + 2
        f_row = fill(C_LIGHT) if i % 2 == 0 else fill(C_WHITE)
        for col, val in enumerate([
            i, stat["파일명"], stat["상대 경로"],
            len(stat["슬라이드 수"]), stat["총 횟수"]
        ], 1):
            c = ws2.cell(row=row, column=col,
                         value=sanitize(val) if isinstance(val, str) else val)
            c.font = DATA_FONT
            c.alignment = CENTER if col in (1, 4, 5) else LEFT
            c.border = thin_border()
            c.fill = f_row
        ws2.row_dimensions[row].height = 18

    # 합계 행
    total_row = len(file_stats) + 3
    ws2.cell(row=total_row, column=3, value="합  계").font = Font(
        name="Arial", bold=True, size=10)
    ws2.cell(row=total_row, column=3).alignment = CENTER
    ws2.cell(row=total_row, column=3).fill = fill(C_MID)

    for col in (4, 5):
        c = ws2.cell(row=total_row, column=col,
                     value=f"=SUM({get_column_letter(col)}3:"
                           f"{get_column_letter(col)}{total_row-1})")
        c.font = Font(name="Arial", bold=True, size=10, color=C_WHITE)
        c.alignment = CENTER
        c.fill = fill(C_MID)
        c.border = thin_border()
    ws2.row_dimensions[total_row].height = 20

    wb.save(output_path)
    print(f"✅ 엑셀 저장 완료: {output_path}")


# ── 인자 파싱 ─────────────────────────────────────────────────────────────────

def parse_args(argv: list[str]) -> tuple[Path, str, Path]:
    """
    공백 포함 경로를 처리합니다.
    규칙: 마지막 인자가 .xlsx → 출력 파일
          그 직전 인자 → 검색어
          나머지 → 경로 (공백 재결합)
    """
    if len(argv) < 2:
        print(__doc__)
        sys.exit(1)

    # 출력 파일 분리
    if argv[-1].lower().endswith(".xlsx"):
        output_arg = argv[-1]
        rest = argv[:-1]
    else:
        output_arg = None
        rest = argv

    if len(rest) < 2:
        print("오류: 디렉토리 경로와 검색어를 모두 입력하세요.")
        sys.exit(1)

    keyword_raw = rest[-1]                          # 마지막 인자 = 검색어(들)
    keywords    = [k.strip() for k in keyword_raw.split(",") if k.strip()]
    dir_parts   = rest[:-1]                         # 나머지 = 경로 조각

    raw_path = " ".join(dir_parts)
    root_dir = Path(raw_path).expanduser().resolve()

    if not root_dir.is_dir():
        print(f"오류: '{root_dir}' 는 유효한 디렉토리가 아닙니다.")
        print('경로에 공백이 있다면 따옴표로 감싸세요:')
        print('  python pptx_search.py "C:\\Users\\내 문서\\강의" 파인튜닝')
        sys.exit(1)

    if output_arg:
        output_path = Path(output_arg)
    else:
        output_path = Path("검색_결과.xlsx")

    return root_dir, keywords, output_path


# ── 진입점 ───────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)

    root_dir, keywords, output_path = parse_args(sys.argv[1:])

    print(f"📂 탐색 경로 : {root_dir}")
    print(f"🔍 검색어    : {', '.join(keywords)}\n")

    results = search_pptx_files(root_dir, keywords)

    if results:
        kw_label = ", ".join(keywords)
        save_to_excel(results, kw_label, output_path)
        import json as _json, os as _os
        _tmp_path = _os.path.join(_os.getcwd(), '_results_tmp.json')
        open(_tmp_path, 'w', encoding='utf-8').write(_json.dumps(results, ensure_ascii=False))
    else:
        print(f'\n검색어 "{keyword}" 가 포함된 슬라이드를 찾지 못했습니다.')


if __name__ == "__main__":
    main()
