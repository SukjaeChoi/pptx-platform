"""
text_analyzer.py
────────────────────────────────────────────────────────────────────
pptx-platform 텍스트 분석 탭 백엔드 (CLI 스크립트)

사용법:
  python text_analyzer.py <디렉토리>
  python text_analyzer.py --files "file1.pptx" "file2.pptx" ...

분석 내용:
  · 형태소 분석 (rhinoMorph)
  · 단어 빈도 (명사 Top30, 내용어 Top20)
  · 품사 분포
  · 명사 공출현 네트워크 (슬라이드 단위)
  · TF-IDF 핵심 키워드 (파일별)

결과: _analysis_tmp.json (SCRIPTS_DIR 저장 후 서버가 읽어감)
"""

import sys
import os
import re
import json
import math
import argparse
from pathlib import Path
from collections import Counter
from itertools import combinations
from datetime import datetime

try:
    from pptx import Presentation
except ImportError:
    print("필수 패키지: pip install python-pptx")
    sys.exit(1)

# ── rhinoMorph 초기화 ────────────────────────────────────────────────────────
try:
    import rhinoMorph
    _rn = rhinoMorph.startRhino()
    print("[OK] rhinoMorph 초기화 완료")
except ImportError:
    print("[ERROR] rhinoMorph 모듈을 찾을 수 없습니다.")
    print("        현재 Python 환경에 rhinoMorph가 설치되어 있는지 확인하세요.")
    print("        server.js의 PYTHON 변수가 rhinoMorph가 설치된 Python을 가리키는지 확인하세요.")
    sys.exit(1)
except Exception as e:
    print(f"[ERROR] rhinoMorph 초기화 실패: {e}")
    sys.exit(1)

NOUN_POS    = ['NNG', 'NNP', 'SL', 'XR']  # 일반명사, 고유명사, 외래어, 어근
VERB_POS    = ['VV']
ADJ_POS     = ['VA']                        # XR은 NOUN_POS로 이동
ADV_POS     = ['MAG', 'MAJ']
CONTENT_POS = NOUN_POS + VERB_POS + ADJ_POS + ADV_POS

def is_valid(w):
    return len(w) > 1 and not re.fullmatch(r'[\d\W]+', w)

def extract_slides(filepath):
    prs = Presentation(filepath)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        slides.append({'num': i + 1, 'text': ' '.join(texts)})
    return slides

def morph_with_pos(text):
    """전체 형태소+품사 분석."""
    if not text.strip():
        return [], []
    try:
        return rhinoMorph.wholeResult_list(_rn, text)
    except Exception as e:
        print(f"\n  [WARN] 형태소 분석 오류: {e}")
        return [], []

def morph_nouns(text):
    """명사(NNG, NNP)만 추출."""
    if not text.strip():
        return []
    try:
        words = rhinoMorph.onlyMorph_list(_rn, text, pos=NOUN_POS)
        return [w for w in words if is_valid(w)]
    except Exception as e:
        print(f"\n  [WARN] 형태소 분석 오류: {e}")
        return []

def compute_tfidf(doc_lists):
    N = len(doc_lists)
    df = Counter()
    for words in doc_lists:
        for w in set(words):
            df[w] += 1
    results = []
    for words in doc_lists:
        total = len(words) or 1
        tf = Counter(words)
        results.append({
            w: round((c / total) * (math.log((N + 1) / (df[w] + 1)) + 1), 5)
            for w, c in tf.items()
        })
    return results

def analyze(pptx_files, scripts_dir):
    n = len(pptx_files)
    print(f"\n📊 텍스트 분석 시작 — {n}개 파일\n{'─' * 60}")

    all_nouns = []
    all_content = []
    all_pos_cnt = Counter()
    cooccurrence = Counter()
    doc_nouns = []
    file_results = []

    for idx, fpath in enumerate(pptx_files, 1):
        fname = Path(fpath).name
        print(f"[{idx:>3}/{n}] {fname}", end=" ... ", flush=True)
        try:
            slides = extract_slides(fpath)
            full_text = ' '.join(s['text'] for s in slides)

            morphs, poses = morph_with_pos(full_text)
            all_pos_cnt += Counter(poses)

            nouns   = [m for m, p in zip(morphs, poses) if p in NOUN_POS    and is_valid(m)]
            content = [m for m, p in zip(morphs, poses) if p in CONTENT_POS and is_valid(m)]

            all_nouns.extend(nouns)
            all_content.extend(content)
            doc_nouns.append(nouns)

            for slide in slides:
                snoun = morph_nouns(slide['text'])
                unique = list(set(snoun))
                for a, b in combinations(unique, 2):
                    cooccurrence[tuple(sorted([a, b]))] += 1

            top4 = [w for w, _ in Counter(nouns).most_common(4)]
            print(f"슬라이드 {len(slides)}개  명사 {len(nouns):,}개  [{', '.join(top4)}]")

            file_results.append({
                'filename': fname, 'fullpath': str(fpath),
                'slides': len(slides), 'noun_count': len(nouns),
                'top_nouns': [w for w, _ in Counter(nouns).most_common(8)],
            })
        except Exception as e:
            print(f"오류: {e}")
            file_results.append({'filename': fname, 'fullpath': str(fpath), 'error': str(e)})
            doc_nouns.append([])

    if doc_nouns:
        tfidf_list = compute_tfidf(doc_nouns)
        for fr, scores in zip(file_results, tfidf_list):
            if 'error' not in fr:
                top = sorted(scores.items(), key=lambda x: -x[1])[:12]
                fr['top_tfidf'] = [{'word': w, 'score': s} for w, s in top if s > 0]

    pos_groups = {
        '일반명사': all_pos_cnt.get('NNG', 0),
        '고유명사': all_pos_cnt.get('NNP', 0),
        '외래어':   all_pos_cnt.get('SL',  0),
        '어근':     all_pos_cnt.get('XR',  0),
        '동사':     all_pos_cnt.get('VV',  0),
        '형용사':   all_pos_cnt.get('VA',  0),
        '부사':     all_pos_cnt.get('MAG', 0) + all_pos_cnt.get('MAJ', 0),
        '관형사':   all_pos_cnt.get('MM',  0),
        '조사':     sum(all_pos_cnt.get(t, 0) for t in ['JX','JKS','JKB','JKO','JKC','JKG']),
        '어미':     sum(all_pos_cnt.get(t, 0) for t in ['EF','EC','ETN','ETM','EP']),
    }

    top_cooc = cooccurrence.most_common(40)
    cooc_node_set = set()
    for (a, b), _ in top_cooc:
        cooc_node_set.add(a); cooc_node_set.add(b)
    noun_freq_all = Counter(all_nouns)
    cooc_nodes = [{'id': nd, 'freq': noun_freq_all.get(nd, 1)} for nd in cooc_node_set]

    result = {
        'generated':          datetime.now().isoformat(),
        'total_files':        len(pptx_files),
        'total_noun_tokens':  sum(f.get('noun_count', 0) for f in file_results),
        'total_vocab':        len(noun_freq_all),
        'rhino_ok':           True,
        'noun_freq':    [{'word': w, 'count': c} for w, c in noun_freq_all.most_common(30)],
        'wc_freq':      [{'word': w, 'count': c} for w, c in noun_freq_all.most_common(100)],
        'content_freq': [{'word': w, 'count': c} for w, c in Counter(all_content).most_common(20)],
        'pos_groups':   pos_groups,
        'cooccurrence': [{'source': a, 'target': b, 'count': c} for (a, b), c in top_cooc],
        'cooc_nodes':   cooc_nodes,
        'files':        file_results,
    }

    out_path = os.path.join(scripts_dir, '_analysis_tmp.json')
    with open(out_path, 'w', encoding='utf-8') as f:
        json.dump(result, f, ensure_ascii=False)

    ok = sum(1 for fr in file_results if 'error' not in fr)
    total_nouns = sum(f.get('noun_count', 0) for f in file_results)
    print(f"\n{'─' * 60}")
    print(f"✅ 분석 완료 — 성공 {ok}/{n}개 | 명사 토큰 {total_nouns:,}개 | 어휘 {len(noun_freq_all):,}종")
    print(f"   결과: {out_path}")

def main():
    parser = argparse.ArgumentParser(description='PPTX 텍스트 분석기')
    parser.add_argument('directory', nargs='?', help='분석할 디렉토리 경로')
    parser.add_argument('--files', nargs='+', metavar='FILE', help='분석할 특정 PPTX 파일 목록')
    args = parser.parse_args()

    scripts_dir = os.path.dirname(os.path.abspath(__file__))

    if args.files:
        pptx_files = [f for f in args.files if f.lower().endswith('.pptx') and os.path.isfile(f)]
        if not pptx_files:
            print("오류: 유효한 PPTX 파일이 없습니다.")
            sys.exit(1)
        print(f"📂 파일 지정 모드 — {len(pptx_files)}개 파일")
    elif args.directory:
        root = Path(args.directory)
        if not root.is_dir():
            print(f"오류: 디렉토리를 찾을 수 없습니다: {args.directory}")
            sys.exit(1)
        pptx_files = sorted([
            str(p) for p in root.rglob('*.pptx') if not p.name.startswith('~$')
        ])
        if not pptx_files:
            print("오류: PPTX 파일이 없습니다.")
            sys.exit(1)
        print(f"📂 탐색 경로: {root}")
    else:
        parser.print_help()
        sys.exit(1)

    analyze(pptx_files, scripts_dir)

if __name__ == '__main__':
    main()
