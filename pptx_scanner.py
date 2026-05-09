"""
pptx_scanner.py
────────────────
지정한 디렉토리(및 하위 디렉토리)에서 모든 .pptx 파일을 탐색하여
파일 제목과 슬라이드 수를 엑셀 파일로 저장합니다.

사용법:
    python pptx_scanner.py <디렉토리_경로> [출력_파일명.xlsx]

예시:
    python pptx_scanner.py /home/user/presentations
    python pptx_scanner.py ./docs report.xlsx
"""

import sys
import os
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
except ImportError:
    print("필수 패키지가 없습니다. 다음 명령어로 설치하세요:")
    print("  pip install python-pptx openpyxl")
    sys.exit(1)

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
except ImportError:
    print("필수 패키지가 없습니다. 다음 명령어로 설치하세요:")
    print("  pip install openpyxl")
    sys.exit(1)


# ── PPTX 파싱 ────────────────────────────────────────────────────────────────

import re

def sanitize(text: str) -> str:
    """openpyxl이 허용하지 않는 제어 문자를 제거하고 공백을 정리합니다."""
    if not isinstance(text, str):
        return text
    # 줄바꿈·탭 → 공백
    text = text.replace("\r\n", " ").replace("\r", " ").replace("\n", " ").replace("\t", " ")
    # XML 1.0 불허 제어 문자 제거 (U+0000-U+0008, U+000B-U+000C, U+000E-U+001F, U+007F-U+009F)
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", text)
    # 연속 공백 정리
    return re.sub(r" +", " ", text).strip()


def get_pptx_title(prs: Presentation) -> str:
    """프레젠테이션 제목을 추출합니다 (여러 방법 순차 시도)."""
    # 1) Core Properties (파일 메타데이터)
    try:
        title = prs.core_properties.title
        if title and title.strip():
            return sanitize(title)
    except Exception:
        pass

    # 2) 첫 번째 슬라이드의 제목 플레이스홀더
    if prs.slides:
        first_slide = prs.slides[0]
        # placeholder type 15 = TITLE, 13 = CENTER_TITLE
        for ph in first_slide.placeholders:
            if ph.placeholder_format.idx in (0, 1):   # 0: title, 1: body (fallback)
                try:
                    text = ph.text_frame.text.strip()
                    if text:
                        return sanitize(text)
                except Exception:
                    pass

        # 3) 첫 번째 슬라이드 텍스트 박스 중 가장 첫 번째 비어있지 않은 텍스트
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    return sanitize(text)

    return "(제목 없음)"


def scan_pptx_files(root_dir: Path) -> list[dict]:
    """root_dir 이하 모든 .pptx 파일을 탐색하여 정보를 반환합니다."""
    results = []
    pptx_files = sorted(root_dir.rglob("*.pptx"))

    if not pptx_files:
        print(f"⚠  '{root_dir}' 에서 .pptx 파일을 찾지 못했습니다.")
        return results

    print(f"총 {len(pptx_files)}개의 .pptx 파일을 발견했습니다.\n")

    for idx, filepath in enumerate(pptx_files, 1):
        rel_path = filepath.relative_to(root_dir)
        print(f"[{idx:>3}/{len(pptx_files)}] {rel_path}", end=" ... ", flush=True)

        try:
            prs = Presentation(filepath)
            title = get_pptx_title(prs)
            slide_count = len(prs.slides)
            status = "OK"
        except PackageNotFoundError:
            title = "(손상된 파일)"
            slide_count = None
            status = "오류: 파일 손상"
        except Exception as e:
            title = "(읽기 실패)"
            slide_count = None
            status = f"오류: {e}"

        results.append({
            "번호":       idx,
            "파일명":     filepath.name,
            "상대 경로":  str(rel_path),
            "절대 경로":  str(filepath),
            "제목":       title,
            "슬라이드 수": slide_count,
            "상태":       status,
        })

        if status == "OK":
            print(f"슬라이드 {slide_count}개  |  {title[:40]}")
        else:
            print(status)

    return results


# ── 엑셀 저장 ─────────────────────────────────────────────────────────────────

HEADER_FILL  = PatternFill("solid", start_color="1F3864")   # 진한 파란색
HEADER_FONT  = Font(name="Arial", bold=True, color="FFFFFF", size=11)
DATA_FONT    = Font(name="Arial", size=10)
ALT_FILL     = PatternFill("solid", start_color="EBF0FA")   # 연한 파란색
CENTER       = Alignment(horizontal="center", vertical="center", wrap_text=True)
LEFT         = Alignment(horizontal="left",   vertical="center", wrap_text=True)
THIN_BORDER  = Border(
    left=Side(style="thin", color="D0D0D0"),
    right=Side(style="thin", color="D0D0D0"),
    top=Side(style="thin", color="D0D0D0"),
    bottom=Side(style="thin", color="D0D0D0"),
)


def save_to_excel(results: list[dict], output_path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "PPTX 목록"

    # ── 상단 정보 ────────────────────────────────────────────────────────────
    ws.merge_cells("A1:G1")
    ws["A1"] = "PPTX 파일 목록"
    ws["A1"].font = Font(name="Arial", bold=True, size=14, color="1F3864")
    ws["A1"].alignment = CENTER

    ws.merge_cells("A2:G2")
    ws["A2"] = f"생성일시: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}  |  총 {len(results)}개 파일"
    ws["A2"].font = Font(name="Arial", size=9, color="666666")
    ws["A2"].alignment = CENTER

    ws.row_dimensions[1].height = 28
    ws.row_dimensions[2].height = 16

    # ── 헤더 ─────────────────────────────────────────────────────────────────
    HEADERS = ["번호", "파일명", "상대 경로", "절대 경로", "제목", "슬라이드 수", "상태"]
    COL_WIDTHS = [7, 30, 40, 55, 40, 13, 12]

    for col, (header, width) in enumerate(zip(HEADERS, COL_WIDTHS), 1):
        cell = ws.cell(row=3, column=col, value=header)
        cell.font      = HEADER_FONT
        cell.fill      = HEADER_FILL
        cell.alignment = CENTER
        cell.border    = THIN_BORDER
        ws.column_dimensions[get_column_letter(col)].width = width

    ws.row_dimensions[3].height = 20

    # ── 데이터 ───────────────────────────────────────────────────────────────
    for row_idx, record in enumerate(results, 4):
        fill = ALT_FILL if row_idx % 2 == 0 else PatternFill()

        values = [
            record["번호"],
            record["파일명"],
            record["상대 경로"],
            record["절대 경로"],
            record["제목"],
            record["슬라이드 수"],
            record["상태"],
        ]
        alignments = [CENTER, LEFT, LEFT, LEFT, LEFT, CENTER, CENTER]

        for col, (val, align) in enumerate(zip(values, alignments), 1):
            cell = ws.cell(row=row_idx, column=col, value=sanitize(val) if isinstance(val, str) else val)
            cell.font      = DATA_FONT
            cell.alignment = align
            cell.border    = THIN_BORDER
            if fill.fill_type:
                cell.fill = fill

            # 오류 행 강조
            if record["상태"] != "OK":
                cell.font = Font(name="Arial", size=10, color="C00000")

        ws.row_dimensions[row_idx].height = 18

    # ── 요약 ─────────────────────────────────────────────────────────────────
    summary_row = len(results) + 5
    ws.cell(row=summary_row, column=5, value="총 슬라이드 수 합계").font = Font(
        name="Arial", bold=True, size=10)
    ws.cell(row=summary_row, column=5).alignment = CENTER

    ok_rows = [r for r in results if r["슬라이드 수"] is not None]
    if ok_rows:
        first_data = 4
        last_data  = len(results) + 3
        ws.cell(row=summary_row, column=6,
                value=f"=SUM(F{first_data}:F{last_data})").font = Font(
            name="Arial", bold=True, size=10)
        ws.cell(row=summary_row, column=6).alignment = CENTER
        ws.cell(row=summary_row, column=6).fill = PatternFill(
            "solid", start_color="D9E1F2")

    # ── 틀 고정 (헤더) ────────────────────────────────────────────────────────
    ws.freeze_panes = "A4"

    wb.save(output_path)
    print(f"\n✅ 엑셀 파일 저장 완료: {output_path}")


# ── 진입점 ───────────────────────────────────────────────────────────────────

def parse_args(argv: list[str]) -> tuple[Path, Path]:
    """
    공백이 포함된 경로를 올바르게 처리합니다.
    - 마지막 인자가 .xlsx로 끝나면 출력 파일로, 나머지를 경로로 결합합니다.
    - 그렇지 않으면 모든 인자를 경로로 결합합니다.
    """
    if not argv:
        print(__doc__)
        sys.exit(1)

    # 마지막 인자가 출력 파일(.xlsx)인지 확인
    if argv[-1].lower().endswith(".xlsx"):
        dir_parts  = argv[:-1]
        output_arg = argv[-1]
    else:
        dir_parts  = argv
        output_arg = None

    # 공백 포함 경로 재결합: 인자들을 공백으로 이어 붙임
    raw_path = " ".join(dir_parts)
    root_dir = Path(raw_path).expanduser().resolve()

    if not root_dir.is_dir():
        # 재결합 후에도 실패하면 명확한 메시지 출력
        print(f"오류: '{root_dir}' 는 유효한 디렉토리가 아닙니다.")
        print("경로에 공백이 있을 경우 따옴표로 감싸는 것을 권장합니다.")
        print('  예시: python pptx_scanner.py "C:\\Users\\내 문서\\강의"')
        sys.exit(1)

    if output_arg:
        output_path = Path(output_arg)
    else:
        output_path = Path("pptx_목록.xlsx")

    return root_dir, output_path


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    root_dir, output_path = parse_args(sys.argv[1:])

    print(f"📂 탐색 경로: {root_dir}\n{'─' * 60}")
    results = scan_pptx_files(root_dir)

    if results:
        save_to_excel(results, output_path)
        import json as _json, os as _os
        _tmp_path = _os.path.join(_os.getcwd(), '_results_tmp.json')
        open(_tmp_path, 'w', encoding='utf-8').write(_json.dumps(results, ensure_ascii=False))
    else:
        print("저장할 데이터가 없습니다.")


if __name__ == "__main__":
    main()
