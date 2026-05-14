import sys
import os
import json
import argparse

# 분석 유형별 파일별 프롬프트 (3문장 이하 요약 기준)
PER_FILE_PROMPTS = {
    'summary':  '다음 문서의 핵심 내용을 3문장 이하로 간결하게 요약해주세요:',
    'topics':   '다음 문서의 핵심 주제와 키워드를 3문장 이하로 간략히 정리해주세요:',
    'qa':       ('다음 문서를 읽고 반드시 아래 형식으로만 출력하세요. 다른 설명은 쓰지 마세요.\n'
                 'Q: [핵심 질문]\nA: [간결한 답변]\nQ: [핵심 질문]\nA: [간결한 답변]'),
    'compare':  '다음 문서의 주요 내용과 특징을 3문장 이하로 간략히 정리해주세요:',
    'keywords': '다음 문서의 핵심 키워드 3~5개를 3문장 이하로 설명해주세요:',
}

# 분석 유형별 종합 프롬프트
OVERALL_PROMPTS = {
    'summary':  '다음은 각 문서의 개별 요약입니다. 전체 내용을 종합하여 최종 요약을 작성해주세요:',
    'topics':   '다음은 각 문서의 개별 주제 분석입니다. 전체 문서의 공통 주제와 주제 분포를 종합 정리해주세요:',
    'qa':       ('다음은 각 문서의 개별 Q&A입니다. '
                 '전체 문서를 아우르는 종합 Q&A를 Q:/A: 형식으로 작성해주세요:'),
    'compare':  '다음은 각 문서의 개별 요약입니다. 문서들을 비교 분석하여 공통점, 차이점, 각 문서의 특징을 정리해주세요:',
    'keywords': '다음은 각 문서의 개별 키워드 분석입니다. 전체 문서의 공통 핵심 키워드와 각 문서의 특징적 키워드를 종합 정리해주세요:',
    'custom':   '다음은 각 문서에 대한 개별 분석 결과입니다. 전체 내용을 종합 정리해주세요:',
}


def extract_text_from_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    parts.append(shape.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                parts.append(cell.text.strip())
                if hasattr(shape, 'notes') and shape.notes:
                    notes_text = getattr(shape.notes, 'text', '')
                    if notes_text.strip():
                        parts.append(f'[노트] {notes_text.strip()}')
            if parts:
                slides.append(f'[슬라이드 {i}]\n' + '\n'.join(parts))
        return '\n\n'.join(slides)
    except Exception as e:
        return f'[파일 읽기 오류: {e}]'


def run_ollama(ollama_module, model, prompt, text, filename):
    """단일 파일에 대한 ollama 스트리밍 분석"""
    full_prompt = f'{prompt}\n\n다음은 분석할 문서 내용입니다:\n\n=== {filename} ===\n{text}'
    result = ''
    stream = ollama_module.chat(
        model=model,
        messages=[{'role': 'user', 'content': full_prompt}],
        stream=True,
    )
    for chunk in stream:
        content = chunk['message']['content']
        result += content
        print(content, end='', flush=True)
    return result


def run_ollama_text(ollama_module, model, prompt):
    """미리 구성된 텍스트 그대로 ollama에 전송"""
    result = ''
    stream = ollama_module.chat(
        model=model,
        messages=[{'role': 'user', 'content': prompt}],
        stream=True,
    )
    for chunk in stream:
        content = chunk['message']['content']
        result += content
        print(content, end='', flush=True)
    return result


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--files', nargs='+', default=[])
    parser.add_argument('--dir', default='')
    parser.add_argument('--prompt', default='')
    parser.add_argument('--type', default='summary')
    parser.add_argument('--model', default='phi4-mini')
    args = parser.parse_args()

    analysis_type = args.type

    # 파일 목록 결정
    files = args.files
    if not files and args.dir:
        import glob
        pattern = os.path.join(args.dir, '**', '*.pptx')
        files = [f for f in glob.glob(pattern, recursive=True)
                 if not os.path.basename(f).startswith('~$')]

    if not files:
        print('오류: 분석할 PPTX 파일이 없습니다.', flush=True)
        sys.exit(1)

    try:
        import ollama as ollama_module
    except ImportError:
        print('오류: ollama 패키지가 설치되지 않았습니다.', flush=True)
        print('  pip install ollama 로 설치하세요.', flush=True)
        sys.exit(1)

    # 프롬프트 결정
    if analysis_type == 'custom':
        per_file_prompt = args.prompt or '다음 문서를 분석해주세요:'
        overall_prompt  = (f'다음은 각 문서에 대한 개별 분석 결과입니다. '
                           f'전체 내용을 종합 정리해주세요. (원래 질의: {args.prompt})')
    else:
        per_file_prompt = PER_FILE_PROMPTS.get(analysis_type, args.prompt)
        overall_prompt  = OVERALL_PROMPTS.get(analysis_type, OVERALL_PROMPTS['custom'])

    print(f'📂 {len(files)}개 파일 텍스트 추출 중...', flush=True)

    # 파일별 텍스트 추출
    file_texts = []
    for fp in files:
        fname = os.path.basename(fp)
        print(f'  📄 읽는 중: {fname}', flush=True)
        text = extract_text_from_pptx(fp)
        if len(text) > 3000:
            text = text[:3000] + '\n[... 이하 생략 ...]'
        file_texts.append({'filename': fname, 'path': fp, 'text': text})

    per_file_results = []
    overall = None

    # ── 파일 1개: 단일 분석 ──────────────────────────────────────────────────
    if len(files) == 1:
        f = file_texts[0]
        print(f'\n🤖 [{f["filename"]}] 분석 중 (모델: {args.model})...', flush=True)
        print('─' * 50, flush=True)
        try:
            result = run_ollama(ollama_module, args.model, per_file_prompt, f['text'], f['filename'])
            per_file_results.append({'filename': f['filename'], 'path': f['path'], 'result': result})
        except Exception as e:
            print(f'\n오류: {e}', flush=True)
            sys.exit(1)
        print('\n' + '─' * 50, flush=True)
        print('✅ 분석 완료', flush=True)

    # ── 파일 복수: 파일별 → 종합 ─────────────────────────────────────────────
    else:
        for f in file_texts:
            print(f'\n🤖 [{f["filename"]}] 개별 분석 중 (모델: {args.model})...', flush=True)
            print('─' * 50, flush=True)
            try:
                result = run_ollama(ollama_module, args.model, per_file_prompt, f['text'], f['filename'])
                per_file_results.append({'filename': f['filename'], 'path': f['path'], 'result': result})
            except Exception as e:
                print(f'\n오류: {e}', flush=True)
                per_file_results.append({'filename': f['filename'], 'path': f['path'], 'result': f'[오류: {e}]'})
            print('\n' + '─' * 50, flush=True)

        # 종합 분석: 파일별 결과를 입력으로 사용
        print(f'\n🔗 전체 종합 분석 중...', flush=True)
        print('─' * 50, flush=True)

        combined = '\n\n'.join(
            f'=== {r["filename"]} ===\n{r["result"]}' for r in per_file_results
        )
        if len(combined) > 6000:
            combined = combined[:6000] + '\n[... 이하 생략 ...]'

        overall_full_prompt = f'{overall_prompt}\n\n{combined}'

        try:
            overall = run_ollama_text(ollama_module, args.model, overall_full_prompt)
        except Exception as e:
            print(f'\n오류: {e}', flush=True)
            overall = f'[종합 분석 오류: {e}]'

        print('\n' + '─' * 50, flush=True)
        print('✅ 전체 분석 완료', flush=True)

    # 결과 저장
    tmp_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '_ai_tmp.json')
    with open(tmp_path, 'w', encoding='utf-8') as fout:
        json.dump({
            'per_file': per_file_results,
            'overall':  overall,
            'files':    [os.path.basename(fp) for fp in files],
            'prompt':   args.prompt,
            'type':     analysis_type,
            'model':    args.model,
        }, fout, ensure_ascii=False, indent=2)


if __name__ == '__main__':
    main()
